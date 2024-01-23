import base64
import datetime
import io
import json
import ntpath
import os
import uuid
from os.path import dirname, abspath
from typing import List

import pythoncom
import win32com.client
from pptx import Presentation
from source_engine.chart_source import ChartSource
from sqlalchemy.orm import Session

from charting import ppt_base_path
from charting.model.chart import ChartModel


class Ppt:

    def __init__(self, template: str = 'dr-template.pptm'):
        self.parent_dir = dirname(dirname(abspath(__file__)))
        self.prs = Presentation(pptx=f'{self.parent_dir}/templates/{template}')
        self.db: ChartSource = ChartSource()

    def create(self, data: List[dict], title: str = None, subtitle: str = None,
               suptitle: str = datetime.datetime.today().strftime("%d.%m.%Y")) -> str:

        title = title if title != "" else "Charts"
        subtitle = subtitle if subtitle != "" else None

        self.__add_title_slide(title=title, subtitle=subtitle, suptitle=suptitle)
        self.__add_slides(data=data)
        self.__add_disclaimer()
        return self.__save()

    def __add_title_slide(self, title: str, subtitle: str, suptitle: str):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        title_obj = slide.placeholders[0]
        title_frame = title_obj.text_frame
        title_frame.text = title

        suptitle_obj = slide.placeholders[14]
        suptitle_frame = suptitle_obj.text_frame
        suptitle_frame.text = suptitle

        subtitle_obj = slide.placeholders[15]
        if subtitle is not None:
            subtitle_frame = subtitle_obj.text_frame
            subtitle_frame.text = subtitle
        else:
            sp = subtitle_obj.element
            sp.getparent().remove(sp)

    def __add_slides(self, data: List[dict]):
        with Session(bind=self.db.engine) as session:
            for item in data:
                chart = session.query(ChartModel).get(item['id'])

                slide_layout = self.prs.slide_layouts[3]
                slide = self.prs.slides.add_slide(slide_layout)

                title = item['heading'] if item['heading'] != '' else ', '.join(chart.category.split(','))
                subtitle = item['subtitle'] if item['subtitle'] != '' else ', '.join(chart.region.split(','))
                slide.placeholders[0].text = title
                slide.placeholders[13].text = subtitle

                image_data = base64.b64decode(chart.base64)
                image_stream = io.BytesIO(image_data)
                slide.placeholders[19].insert_picture(image_stream)

    def __add_disclaimer(self):
        slide_layout = self.prs.slide_layouts[17]
        slide = self.prs.slides.add_slide(slide_layout)

        note = slide.placeholders[10]
        sp = note.element
        sp.getparent().remove(sp)

    def get_layout(self):
        for slide in self.prs.slide_layouts:
            for shape in slide.placeholders:
                print('%d %d %s' % (self.prs.slide_layouts.index(slide), shape.placeholder_format.idx, shape.name))

    def __save(self) -> str:
        try:
            path = os.path.join(ppt_base_path, f'{uuid.uuid4().__str__()}.ppt')
            self.prs.save(path)
            filename = ntpath.basename(path)

            pythoncom.CoInitialize()
            powerpoint = win32com.client.gencache.EnsureDispatch('PowerPoint.Application')
            powerpoint.Visible = True
            presentation = powerpoint.Presentations.Open(path)
            presentation.Application.Run(f"{filename}!Modul1.AdjustShapeWidthToFitText")

            presentation.Save()
            presentation.Close()

            powerpoint.Quit()
            pythoncom.CoUninitialize()
        except Exception as e:
            print(e)
            return str(e)

        return path
